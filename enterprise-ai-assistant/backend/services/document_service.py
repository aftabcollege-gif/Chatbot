"""Text extraction from supported document formats."""
from __future__ import annotations

import csv
import io
import json
from pathlib import Path
from typing import Dict, List, Tuple

from bs4 import BeautifulSoup


class ExtractionError(RuntimeError):
    pass


def extract_pdf(path: Path) -> Tuple[str, int]:
    import pdfplumber

    text_parts: List[str] = []
    pages = 0
    with pdfplumber.open(str(path)) as pdf:
        pages = len(pdf.pages)
        for i, page in enumerate(pdf.pages, start=1):
            txt = page.extract_text() or ""
            if txt.strip():
                text_parts.append(f"\n[[page:{i}]]\n{txt}")
    return "\n".join(text_parts), pages


def extract_docx(path: Path) -> Tuple[str, int]:
    import docx

    doc = docx.Document(str(path))
    parts: List[str] = []
    for para in doc.paragraphs:
        style = (para.style.name or "") if para.style else ""
        text = para.text.strip()
        if not text:
            continue
        if style.startswith("Heading"):
            level = "".join(ch for ch in style if ch.isdigit()) or "1"
            parts.append(f"{'#' * int(level)} {text}")
        else:
            parts.append(text)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            parts.append(" | ".join(cells))
    return "\n".join(parts), 1


def extract_xlsx(path: Path) -> Tuple[str, int]:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts: List[str] = []
    for ws in wb.worksheets:
        parts.append(f"## {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    wb.close()
    return "\n".join(parts), 1


def extract_pptx(path: Path) -> Tuple[str, int]:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: List[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"## اسلاید {i}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return "\n".join(parts), len(prs.slides)


def extract_html(path: Path) -> Tuple[str, int]:
    soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="ignore"), "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    text = soup.get_text("\n")
    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
    return "\n".join(lines), 1


def extract_markdown(path: Path) -> Tuple[str, int]:
    return path.read_text(encoding="utf-8", errors="ignore"), 1


def extract_text(path: Path) -> Tuple[str, int]:
    return path.read_text(encoding="utf-8", errors="ignore"), 1


def extract_csv(path: Path) -> Tuple[str, int]:
    content = path.read_text(encoding="utf-8", errors="ignore")
    reader = csv.reader(io.StringIO(content))
    parts = [" | ".join(row) for row in reader if row]
    return "\n".join(parts), 1


def extract_json(path: Path) -> Tuple[str, int]:
    data = json.loads(path.read_text(encoding="utf-8", errors="ignore"))
    return json.dumps(data, ensure_ascii=False, indent=2), 1


EXTRACTORS = {
    "pdf": extract_pdf,
    "docx": extract_docx,
    "xlsx": extract_xlsx,
    "pptx": extract_pptx,
    "html": extract_html,
    "htm": extract_html,
    "md": extract_markdown,
    "markdown": extract_markdown,
    "txt": extract_text,
    "csv": extract_csv,
    "json": extract_json,
    "xml": extract_text,
}


def extract(path: Path, file_type: str) -> Tuple[str, int]:
    fn = EXTRACTORS.get(file_type)
    if fn is None:
        raise ExtractionError(f"استخراج متن برای نوع {file_type} پشتیبانی نمی‌شود.")
    try:
        return fn(path)
    except Exception as exc:  # noqa: BLE001
        raise ExtractionError(f"خطا در استخراج متن: {exc}") from exc
